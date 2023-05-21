from dataclasses import dataclass
import re
from io import BytesIO
from base64 import b64encode
import logging

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE


logging.basicConfig(level=logging.DEBUG)


class BlackAndWhiteTheme:
    font = "Calibri"
    font_color = RGBColor(0, 0, 0)  # Black
    font_size = Pt(32)
    background_color = RGBColor(255, 255, 255)  # White


class GreyScaleTheme1:
    font = "Times New Roman"
    font_color = RGBColor(211, 211, 211) 
    font_size = Pt(32)
    background_color = RGBColor(64, 64, 64)


class GreyScaleTheme2:
    font = "Arial"
    font_color = RGBColor(51, 51, 51)  # Dark gray
    font_size = Pt(28)
    background_color = RGBColor(240, 240, 240) # Light gray


class ForestTheme:
    font = "Georgia"
    font_color = RGBColor(51, 102, 0)  # Dark green
    font_size = Pt(30)
    background_color = RGBColor(204, 255, 204)  # Light green


THEME_MAPPING = {
    "blackAndWhiteTheme": BlackAndWhiteTheme,
    "greyScaleTheme1": GreyScaleTheme1,
    "greyScaleTheme2": GreyScaleTheme2,
    "forestTheme": ForestTheme
}


@dataclass
class Slide:
    # class to hold the title and content for each slide in the presentation
    title: str = ""
    content: list[str] or None = None


def get_info_from_line(line, info_keyword):
    # replace 'title' or 'content' with nothing, returning the title or content of the slide
    return re.sub(fr"{info_keyword}:?\s*", "", line, flags=re.IGNORECASE).strip().strip('"').strip("'")


def get_slide_content_from_lines(lines, idx):
    # returns the bullet points for a slide as a list of strings
    slide_content = []

    while idx < len(lines):
        if lines[idx].lower().startswith("slide"):
            break

        if lines[idx].lower().startswith("content"):
            content = get_info_from_line(lines[idx], "content")
            if content:
                slide_content.append(content)
        else:
            slide_content.append(lines[idx].strip("-").strip())

        idx += 1

    return slide_content, idx-1


def text_to_slide_objs(text):
    # returns the slides for the presentation as a list of Slide objects 
    lines = [line for line in text.split("\n") if line]
    slides = []

    presentation_title = get_info_from_line(lines[0], "title")
    slides.append(Slide(presentation_title))
    # 0 index slide count
    slide_count = 0

    line_idx = 1
    while line_idx < len(lines):
        if lines[line_idx].lower().startswith("slide"):
            slides.append(Slide())
            slide_count += 1
        elif lines[line_idx].lower().startswith("title"):
            slides[slide_count].title = get_info_from_line(lines[line_idx], "title")
        elif lines[line_idx].lower().startswith("content"):
            slide_content, line_idx = get_slide_content_from_lines(lines, line_idx)
            slides[slide_count].content = slide_content
        line_idx += 1

    return slides


def font_edit(font, theme, is_body_text):
    if is_body_text:
        font.size = theme.font_size
    font.name = theme.font
    font.color.rgb = theme.font_color


def gen_presentation(slides, theme_str):
    # generates the presentation and returns a base64 enocded bytestring representing the pptx file

    theme = THEME_MAPPING.get(theme_str)
    if not theme:
        logging.debug("Invalid theme received. Using BlackAndWhiteTheme as default.")
        theme = BlackAndWhiteTheme

    # setting up the presentation
    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    bullet_slide_layout = presentation.slide_layouts[1]

    # setting up the title slide
    slide = presentation.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = slides[0].title
    font_edit(title_shape.text_frame.paragraphs[0].font, theme, False)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = theme.background_color

    # setting up the other slides
    for slide_info in slides[1:]:
        slide = presentation.slides.add_slide(bullet_slide_layout)
        # slide.placeholders[1].left = Pt(36.0)  # resets the slide body's left position to the default
        # slide.placeholders[1].top = Pt(150.0)  # sets the slide body's top position 24 points below the default, which is 126 points
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = theme.background_color
        title_shape = slide.shapes.title
        title_shape.text = slide_info.title  # adds the slide's title
        font_edit(title_shape.text_frame.paragraphs[0].font, theme, False)

        # adds the bullet points
        bullet_text_frame = slide.shapes.placeholders[1].text_frame
        bullet_text_frame.text = slide_info.content[0]
        font_edit(bullet_text_frame.paragraphs[0].font, theme, True)
        for i in range(1, len(slide_info.content)):
            paragraph = bullet_text_frame.add_paragraph()
            paragraph.text = slide_info.content[i]
            paragraph.level = 0
            paragraph.space_before = Pt(15)
            font_edit(bullet_text_frame.paragraphs[i].font, theme, True)
        bullet_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # finally saves the finished pptx file
    pres_bytes = BytesIO()
    presentation.save(pres_bytes)
    return b64encode(pres_bytes.getvalue()).decode()
