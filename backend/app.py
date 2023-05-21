import os
import re
import logging
from io import BytesIO
from base64 import b64encode

from flask import Flask, request
from flask_cors import CORS
from flask_restful import Api, Resource
import openai

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

from pptx_utils import text_to_slide_objs


openai.api_key = os.getenv("OPENAI_API_KEY")
MODEL = "gpt-3.5-turbo"
TEMPERATURE = 0.5

# defaults if not in request body
TOPIC_COUNT = 5
SLIDE_COUNT = 10
CHARACTER_LIMIT = 300
CUSTOM_FONT = "Calibri"
CUSTOM_FONT_COLOR = (0, 0, 0)
CUSTOM_FONT_SIZE = 32
CUSTOM_BACKGROUND_COLOR = (255, 255, 255)


logging.basicConfig(level=logging.DEBUG)


def get_chat_completion(prompt, model=MODEL, temperature=TEMPERATURE):

    resp = openai.ChatCompletion.create(
        model=model,
        messages=[{
            "role": "user",
            "content": prompt
        }],
        temperature=temperature
    )

    return resp["choices"][0]["message"]["content"]


def font_edit(font, is_body_text):
    if is_body_text:
        font.size = Pt(CUSTOM_FONT_SIZE)
    font.name = CUSTOM_FONT
    font.color.rgb = RGBColor(CUSTOM_FONT_COLOR[0], CUSTOM_FONT_COLOR[1], CUSTOM_FONT_COLOR[2])


class PresentationIdeas(Resource):

    def post(self):
        body = request.get_json()

        previous_topics = body.get("previousTopics")
        if not (topic_count := body.get("topicCount")):
            topic_count = TOPIC_COUNT

        prompt = f"Generate a list of {topic_count} different topic ideas for a PowerPoint presentation. Put each topic on a new line and keep it concise."
        if previous_topics:
            prompt += f" Do not include any of the following topics: {','.join(previous_topics)}"

        resp_content = get_chat_completion(prompt)
        topics = resp_content.split("\n")
        for i in range(len(topics)):
            topics[i] = re.sub("\d+\.?\s*", "", topics[i]).strip('"').strip("-")

        logging.debug(f"Here are the topics: {topics}")
        return {"topics": topics}, 200


class SlideDeck(Resource):

    def post(self):
        body = request.get_json()

        # collecting data from chatgpt
        if not (topic := body.get("topic")):
            return "Presentation topic is required. Please include a 'topic' field in the request body", 400
        
        slide_count = body.get("slide_count")
        character_limit = body.get("character_limit")
        
        prompt = f'''You are an expert on the topic of "{topic}". Generate the content for a PowerPoint presentation on that topic. The following text, enclosed in ", is the format for a complete presentation that is 3 slides long. Parts that you should fill out start with [ and end with ].\n"\nPresentation Title: [catchy title for your presentation]\nSlide 1\nTitle: [catchy title for Slide 1 (cannot be the word 'Introduction')]\nContent:\n[bulleted list of points for Slide 1] (use full but very concise sentences)\nSlide 2\nTitle: [catchy title for Slide 2]\nContent:\n[bulleted list of points for Slide 2] (use full but very concise sentences)\nSlide 3\nTitle: [clever title for Slide 3 (cannot be the word 'Conclusion')]\nContent:\n[bulleted list of points for Slide 3] (use full but very concise sentences)\n"\nUse this format to create exactly {slide_count if slide_count else SLIDE_COUNT} slides, which should include a conclusive final slide. The slide content must not exceed {character_limit if character_limit else CHARACTER_LIMIT}. Do not add anything before or after the presentation.'''
        resp_content = get_chat_completion(prompt)
        logging.debug("This is the response: ", resp_content)
        
        slides = text_to_slide_objs(resp_content)

        # setting up the presentation
        presentation = Presentation()
        title_slide_layout = presentation.slide_layouts[0]
        bullet_slide_layout = presentation.slide_layouts[1]

        # setting up the title slide
        slide = presentation.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = slides[0].title
        font_edit(title_shape.text_frame.paragraphs[0].font, False)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(CUSTOM_BACKGROUND_COLOR[0], CUSTOM_BACKGROUND_COLOR[1], CUSTOM_BACKGROUND_COLOR[2])

        # setting up the other slides
        for slide_info in slides[1:]:
            slide = presentation.slides.add_slide(bullet_slide_layout)
            # slide.placeholders[1].left = Pt(36.0)  # resets the slide body's left position to the default
            # slide.placeholders[1].top = Pt(150.0)  # sets the slide body's top position 24 points below the default, which is 126 points
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(CUSTOM_BACKGROUND_COLOR[0], CUSTOM_BACKGROUND_COLOR[1], CUSTOM_BACKGROUND_COLOR[2])
            title_shape = slide.shapes.title
            title_shape.text = slide_info.title  # adds the slide's title
            font_edit(title_shape.text_frame.paragraphs[0].font, False)

            # adds the bullet points
            bullet_text_frame = slide.shapes.placeholders[1].text_frame
            bullet_text_frame.text = slide_info.content[0]
            font_edit(bullet_text_frame.paragraphs[0].font, True)
            for i in range(1, len(slide_info.content)):
                paragraph = bullet_text_frame.add_paragraph()
                paragraph.text = slide_info.content[i]
                paragraph.level = 0
                paragraph.space_before = Pt(15)
                font_edit(bullet_text_frame.paragraphs[i].font, True)
            bullet_text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # finally saves the finished pptx file
        pres_bytes = BytesIO()
        presentation.save(pres_bytes)
        pres_bytestring = b64encode(pres_bytes.getvalue()).decode()

        return {"byteString": pres_bytestring}, 200


def create_app():
    app = Flask(__name__)
    CORS(app, origins=os.environ["ALLOWED_HOSTS"])
    app.config["SECRET_KEY"] = os.environ["SECRET_KEY"]
    return app


app = create_app()
api = Api(app)
api.add_resource(PresentationIdeas, "/api/topics")
api.add_resource(SlideDeck, "/api/slideDeck")


if __name__ == "__main__":
    app.run()
    